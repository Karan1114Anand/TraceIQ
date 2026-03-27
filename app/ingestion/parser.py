"""
app/ingestion/parser.py

Multi-format document parser with a dispatcher pattern.
Supports: PDF, DOCX, PPTX, CSV/Excel.
"""

from __future__ import annotations

import csv
from pathlib import Path
from typing import Dict, List, Optional

from loguru import logger


# ---------------------------------------------------------------------------
# PDF Parser
# ---------------------------------------------------------------------------

class PDFParser:
    """Parser for PDF documents using pdfplumber + pymupdf."""

    supported_extensions = [".pdf"]

    def parse(self, file_path: Path) -> Dict:
        logger.info(f"Parsing PDF: {file_path.name}")
        try:
            import pdfplumber
            import fitz  # pymupdf

            text_parts: List[str] = []
            tables: List[Dict] = []

            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    try:
                        page_text = page.extract_text(layout=True)
                        if page_text:
                            text_parts.append(f"\n--- Page {page_num} ---\n{page_text}")
                        for t_idx, table in enumerate(page.extract_tables() or []):
                            if table:
                                tables.append(
                                    {"page": page_num, "table_index": t_idx, "data": table}
                                )
                    except Exception as exc:
                        logger.warning(f"Page {page_num} extraction error: {exc}")

            # Metadata via pymupdf
            meta: Dict = {"page_count": 0}
            try:
                doc = fitz.open(file_path)
                raw = doc.metadata or {}
                meta = {
                    "title": raw.get("title", ""),
                    "author": raw.get("author", ""),
                    "page_count": doc.page_count,
                    "creation_date": raw.get("creationDate", ""),
                }
                doc.close()
            except Exception as exc:
                logger.warning(f"Metadata extraction failed: {exc}")

            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": "pdf",
                "text": "\n".join(text_parts),
                "tables": tables,
                "metadata": meta,
                "page_count": meta.get("page_count", 0),
            }
        except Exception as exc:
            logger.error(f"PDF parse error: {exc}")
            raise


# ---------------------------------------------------------------------------
# DOCX Parser
# ---------------------------------------------------------------------------

class DOCXParser:
    """Parser for Microsoft Word (.docx) documents."""

    supported_extensions = [".docx", ".doc"]

    def parse(self, file_path: Path) -> Dict:
        logger.info(f"Parsing DOCX: {file_path.name}")
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts: List[str] = []
            for para in doc.paragraphs:
                if para.text.strip():
                    if para.style and para.style.name and para.style.name.startswith("Heading"):
                        lvl = para.style.name.replace("Heading ", "")
                        hashes = "#" * (int(lvl) if lvl.isdigit() else 1)
                        text_parts.append(f"\n{hashes} {para.text}\n")
                    else:
                        text_parts.append(para.text)

            tables: List[Dict] = []
            for t_idx, table in enumerate(doc.tables):
                try:
                    data = [[c.text.strip() for c in row.cells] for row in table.rows]
                    if data:
                        tables.append({"table_index": t_idx, "data": data})
                except Exception as exc:
                    logger.warning(f"Table {t_idx} extraction error: {exc}")

            try:
                cp = doc.core_properties
                meta = {
                    "title": cp.title or "",
                    "author": cp.author or "",
                    "created": str(cp.created) if cp.created else "",
                    "modified": str(cp.modified) if cp.modified else "",
                }
            except Exception:
                meta = {}

            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": "docx",
                "text": "\n".join(text_parts),
                "tables": tables,
                "metadata": meta,
                "page_count": 0,
            }
        except Exception as exc:
            logger.error(f"DOCX parse error: {exc}")
            raise


# ---------------------------------------------------------------------------
# PPTX Parser
# ---------------------------------------------------------------------------

class PPTXParser:
    """Parser for Microsoft PowerPoint (.pptx) files."""

    supported_extensions = [".pptx", ".ppt"]

    def parse(self, file_path: Path) -> Dict:
        logger.info(f"Parsing PPTX: {file_path.name}")
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_data: List[Dict] = []
            text_parts: List[str] = []

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text: List[str] = []
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                        if shape.has_table:
                            rows = [
                                " | ".join(c.text.strip() for c in row.cells)
                                for row in shape.table.rows
                            ]
                            slide_text.append("\n".join(rows))

                    notes = ""
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()

                    slides_data.append(
                        {"slide_number": slide_idx, "content": "\n".join(slide_text), "notes": notes}
                    )
                    text_parts.append(f"\n--- Slide {slide_idx} ---\n" + "\n".join(slide_text))
                    if notes:
                        text_parts.append(f"Notes: {notes}")
                except Exception as exc:
                    logger.warning(f"Slide {slide_idx} error: {exc}")

            try:
                cp = prs.core_properties
                meta = {"title": cp.title or "", "author": cp.author or ""}
            except Exception:
                meta = {}

            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": "pptx",
                "text": "\n".join(text_parts),
                "slides": slides_data,
                "metadata": meta,
                "page_count": len(prs.slides),
            }
        except Exception as exc:
            logger.error(f"PPTX parse error: {exc}")
            raise


# ---------------------------------------------------------------------------
# CSV / Excel Parser
# ---------------------------------------------------------------------------

class CSVParser:
    """Parser for CSV and Excel files."""

    supported_extensions = [".csv", ".xlsx", ".xls"]

    def parse(self, file_path: Path) -> Dict:
        logger.info(f"Parsing spreadsheet: {file_path.name}")
        try:
            if file_path.suffix.lower() == ".csv":
                return self._parse_csv(file_path)
            return self._parse_excel(file_path)
        except Exception as exc:
            logger.error(f"Spreadsheet parse error: {exc}")
            raise

    def _parse_csv(self, file_path: Path) -> Dict:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
            sample = fh.read(1024)
            fh.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel
            data = list(csv.reader(fh, dialect))
        headers = data[0] if data else []
        rows = data[1:] if len(data) > 1 else []
        return {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "file_type": "csv",
            "text": self._table_to_text(headers, rows),
            "metadata": {"row_count": len(rows), "column_count": len(headers)},
            "page_count": 1,
        }

    def _parse_excel(self, file_path: Path) -> Dict:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts: List[str] = []
        total_rows = 0
        for name in wb.sheetnames:
            ws = wb[name]
            data = [[str(c) if c is not None else "" for c in row] for row in ws.iter_rows(values_only=True)]
            if not data:
                continue
            headers, rows = data[0], data[1:]
            total_rows += len(rows)
            text_parts.append(f"\n--- Sheet: {name} ---\n{self._table_to_text(headers, rows)}")
        wb.close()
        return {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "file_type": "excel",
            "text": "\n".join(text_parts),
            "metadata": {"sheet_count": len(wb.sheetnames), "total_rows": total_rows},
            "page_count": 1,
        }

    @staticmethod
    def _table_to_text(headers: List[str], rows: List[List[str]], max_rows: int = 100) -> str:
        lines = [" | ".join(headers), "-" * 80] if headers else []
        lines += [" | ".join(r) for r in rows[:max_rows]]
        if len(rows) > max_rows:
            lines.append(f"... ({len(rows) - max_rows} more rows)")
        return "\n".join(lines)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

class ParserDispatcher:
    """Routes files to the correct parser based on file extension."""

    def __init__(self) -> None:
        self._parsers = {
            ".pdf": PDFParser(),
            ".docx": DOCXParser(),
            ".doc": DOCXParser(),
            ".pptx": PPTXParser(),
            ".ppt": PPTXParser(),
            ".csv": CSVParser(),
            ".xlsx": CSVParser(),
            ".xls": CSVParser(),
        }
        logger.info(f"ParserDispatcher ready. Supported: {list(self._parsers.keys())}")

    def parse(self, file_path: Path) -> Optional[Dict]:
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        parser = self._parsers.get(file_path.suffix.lower())
        if not parser:
            logger.warning(f"Unsupported file type: {file_path.suffix}")
            return None
        return parser.parse(file_path)

    def is_supported(self, file_path: Path) -> bool:
        return Path(file_path).suffix.lower() in self._parsers
